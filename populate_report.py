#!/usr/bin/env python3
import sys
import json
import os
from pptx import Presentation

def replace_text_in_shape(shape, replacements):
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            for old, new in replacements.items():
                if old in run.text:
                    run.text = run.text.replace(old, new)

def replace_in_presentation(pptx_path, output_path, data):
    prs = Presentation(pptx_path)
    
    replacements = {
        '[COMPANY NAME]': data.get('companyName', ''),
        '[Owner Name]': data.get('ownerName', ''),
        'January 21, 2026': data.get('reportDate', ''),
        '[YOUR EMAIL]': 'dave@profitablebusinesssolutions.com',
        '[YOUR PHONE]': '',
        '$47,000+': data.get('totalDormantRevenue', ''),
        '$47,000': data.get('totalDormantRevenue', ''),
        '1,440 customers': data.get('dormantCount', '') + ' customers',
        '1,440': data.get('dormantCount', ''),
        '3,200': data.get('totalCustomers', ''),
        '$288k': data.get('sunkCostLow', ''),
        '$720k': data.get('sunkCostHigh', ''),
        '$117,822': data.get('reactivationRevenue', ''),
        '$8,400': data.get('month1Revenue', ''),
        '$15,600': data.get('month2Revenue', ''),
        '$23,800': data.get('month3Revenue', ''),
        '$47,800': data.get('reactivationRevenue', ''),
    }

    for slide in prs.slides:
        for shape in slide.shapes:
            replace_text_in_shape(shape, replacements)
            if shape.shape_type == 6:
                for s in shape.shapes:
                    replace_text_in_shape(s, replacements)

    prs.save(output_path)
    print(f"Saved: {output_path}")

if __name__ == '__main__':
    data = json.loads(sys.argv[1])
    output = sys.argv[2] if len(sys.argv) > 2 else '/tmp/report_output.pptx'
    template = os.path.join(os.path.dirname(__file__), 'report_template.pptx')
    replace_in_presentation(template, output, data)
